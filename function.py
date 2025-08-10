# -*- coding: utf-8 -*-
import os

def convert_to_text(filepath: str) -> str:
    ext = os.path.splitext(filepath)[-1].lower()

    if ext == ".pdf":
        import fitz  # PyMuPDF
        text = ""
        with fitz.open(filepath) as doc:
            for page in doc:
                text += page.get_text()
        return text

    elif ext == ".docx":
        from docx import Document
        doc = Document(filepath)
        return "\n".join([para.text for para in doc.paragraphs])

    elif ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(filepath)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text

    elif ext == ".txt":
        return open(filepath, "r", encoding="utf-8").read()
    
    elif ext == ".csv":
        import pandas as pd
        df = pd.read_csv(filepath, encoding="utf-8")
        return df.to_csv(index=False)
    
    elif ext in (".xls", ".xlsx"):
        import pandas as pd

        # 讀取所有工作表；pandas 會回傳 dict: {sheet_name: DataFrame}
        all_sheets = pd.read_excel(
            filepath,
            sheet_name=None,       # 讀全部
            header=None,           # 不把第一列當欄名，避免遺漏
            engine="openpyxl" if ext == ".xlsx" else "xlrd"
        )

        # 把每張工作表轉成文字
        text_parts = []
        for name, df in all_sheets.items():
            # 將 NaN 補空字串、改成 str，再以 tab 串列 → 每列加換行
            body = (
                df.fillna("")
                .astype(str)
                .apply(lambda row: "\t".join(row), axis=1)
                .str.cat(sep="\n")
            )
            text_parts.append(f"--- 工作表: {name} ---\n{body}\n")

        return "\n".join(text_parts)
    
    elif ext == ".json":
        import json

        def _flatten(obj, prefix=""):
            items = {}
            if isinstance(obj, dict):
                for k, v in obj.items():
                    new_key = f"{prefix}.{k}" if prefix else k
                    items.update(_flatten(v, new_key))
            elif isinstance(obj, list):
                for i, v in enumerate(obj):
                    new_key = f"{prefix}[{i}]" if prefix else f"[{i}]"
                    items.update(_flatten(v, new_key))
            else:
                # 對於基本類型，直接存
                items[prefix] = obj
            return items

        text_parts = []
        with open(filepath, "r", encoding="utf-8") as f:
            # 讀前一小段判斷是不是 NDJSON（每行一個 JSON 物件）
            sample = f.read(2048)
            f.seek(0)
            is_ndjson = False
            try:
                lines = sample.strip().splitlines()
                if len(lines) > 1:
                    # 嘗試解析前幾行，如果都能 parse 成 json，認為是 NDJSON
                    for line in lines[:5]:
                        json.loads(line)
                    is_ndjson = True
            except Exception:
                is_ndjson = False

            if is_ndjson:
                for idx, line in enumerate(f, start=1):
                    line = line.strip()
                    if not line:
                        continue
                    try:
                        obj = json.loads(line)
                    except json.JSONDecodeError:
                        continue  # 跳過 parse 失敗行
                    flat = _flatten(obj)
                    body = "\n".join(f"{k}: {v}" for k, v in flat.items())
                    text_parts.append(f"--- RECORD {idx} ---\n{body}")
            else:
                try:
                    data = json.load(f)
                except json.JSONDecodeError:
                    # 如果整個 JSON 讀不進來，退回原始文字
                    f.seek(0)
                    return f.read()
                if isinstance(data, list):
                    for idx, item in enumerate(data, start=1):
                        flat = _flatten(item)
                        body = "\n".join(f"{k}: {v}" for k, v in flat.items())
                        text_parts.append(f"--- ITEM {idx} ---\n{body}")
                elif isinstance(data, dict):
                    flat = _flatten(data)
                    body = "\n".join(f"{k}: {v}" for k, v in flat.items())
                    text_parts.append(body)
                else:
                    text_parts.append(str(data))

        return "\n\n".join(text_parts)
  
    else:
        raise ValueError(f"❌ 不支援的教材格式：{ext}")

def multiline_input(prompt="你（可多行，空行送出）："):
    print(prompt)
    lines = []
    while True:
        line = input()
        if line.strip() == "":
            break
        lines.append(line)
    return "\n".join(lines)
